import os
import sys
from pptx import Presentation

def main():
    template_path = r"c:\Users\raman\Downloads\final pro third rev.pptx"
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print("Failed to load presentation:", e)
        sys.exit(1)

    # Finding generic layouts (0 -> Title, 1 -> Content)
    title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None
    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else title_layout

    if not title_layout or not content_layout:
        print("Template does not have standard slide layouts.")
        sys.exit(1)

    # 1. Project Title with domain and abstract
    slide = prs.slides.add_slide(title_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Smart Crop Weather Intelligence\n(Weather-Based Crop Prediction)"
    if len(slide.placeholders) > 1:
        # Assuming the second placeholder is suitable for subtitle/abstract
        slide.placeholders[1].text = "Domain: Machine Learning & Agricultural Data Analytics\n\nAbstract: Analyzing climate inputs to recommend optimal crops using Decision Trees and bridging meteorology with agriculture."

    # 2. Introduction
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Introduction"
    
    # Locate body placeholder
    content_ph = None
    if len(slide.placeholders) > 1:
        content_ph = slide.placeholders[1]
        
    if content_ph:
        content_ph.text = "Scope & Objectives:\n- Provide accessible, reliable crop recommendations.\n- Maximize yield and mitigate financial risks.\n\nSystem Description:\n- Full-stack app with dynamic web GUI mapped to robust Python ML backend.\n\nProposed System:\n- Responsive search portal with automated pipeline generating 1000+ crops."

    # 3. Overall Description
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Overall Description"
        
    content_ph = None
    if len(slide.placeholders) > 1:
        content_ph = slide.placeholders[1]
        
    if content_ph:
        content_ph.text = "Hardware Requirements:\n- Core CPU, 2GB+ RAM, 500MB+ Storage\n\nSoftware Requirements:\n- Windows/Linux/macOS\n- Python, Web Browser\n- Pandas, Scikit-Learn, HTML/JS/CSS\n\nCharacteristics:\n- Responsive UI, scalable DB\n- Cross-browser compatible"

    # Save logic
    output_path = r"s:\MINI project sample\crop weather\website\final_pro_third_rev_Updated.pptx"
    prs.save(output_path)
    print("Successfully added slides to " + output_path)

if __name__ == '__main__':
    main()
