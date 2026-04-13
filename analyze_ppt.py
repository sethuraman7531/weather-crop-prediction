import traceback
import sys
try:
    from pptx import Presentation
    prs = Presentation(r"c:\Users\raman\Downloads\College PPT Templet.pptx")
    print("Presentation loaded. Number of slides: {}".format(len(prs.slides)))
    for i, slide in enumerate(prs.slides):
        print("Slide {}:".format(i))
        for shape in slide.shapes:
            if shape.has_text_frame:
                print("  - Text: " + repr(shape.text))

    print("\nSlide Layouts:")
    for i, layout in enumerate(prs.slide_layouts):
        print("Layout {}: {}".format(i, layout.name))
        for placeholder in layout.placeholders:
            print("  - Placeholder [id={}] -> {}".format(placeholder.placeholder_format.idx, placeholder.name))
except Exception as e:
    print("Error:")
    traceback.print_exc()
