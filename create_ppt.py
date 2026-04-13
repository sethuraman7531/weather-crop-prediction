# -*- coding: utf-8 -*-
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Could not import python-pptx. Please ensure it is installed.")
    sys.exit(1)

def main():
    try:
        prs = Presentation()
    except Exception as e:
        print("Failed to initialize presentation: " + str(e))
        sys.exit(1)

    # Slide 1: Title
    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Smart Crop Weather Intelligence\n(Weather-Based Crop Prediction)"
    subtitle.text = "Domain: Machine Learning & Agricultural Data Analytics\n\nAbstract: Analyzing climate inputs to recommend optimal crops using Decision Trees and bridging meteorology with agriculture."

    # Slide 2: Intro
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introduction"
    content = slide.placeholders[1]
    content.text = "Scope & Objectives:\nProvide accessible, reliable crop recommendations.\nMaximize yield and mitigate financial risks.\n\nSystem Description:\nFull-stack app with dynamic web GUI mapped to robust Python ML backend.\n\nProposed System:\nResponsive search portal with automated pipeline generating 1000+ crops."

    # Slide 3: Overall Description
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Overall Description"
    content = slide.placeholders[1]
    content.text = "Hardware Requirements:\nCore CPU, 2GB+ RAM, 500MB+ Storage\n\nSoftware Requirements:\nWindows/Linux/macOS\nPython, Web Browser\nPandas, Scikit-Learn, HTML/JS/CSS\n\nCharacteristics:\nResponsive UI, scalable DB\nCross-browser compatible"

    # Slide 4: Specific Requirements
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Specific Requirements"
    content = slide.placeholders[1]
    content.text = "Functional Requirements:\nAccept temp, rain, humidity inputs -> return prediction.\nCatalog search for localized needs.\n\nNon-Functional Requirement:\nScalable to 1000+ crops, fast load speed.\n\nExternal Interfaces:\nParticle.js backgrounds, FontAwesome.\n\nPerformance:\nML execution avoids memory overflow; Lookups < 500ms."

    # Slide 5: Action Plan
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Action Plan (Gantt Overview)"
    content = slide.placeholders[1]
    content.text = "T1 [Week 1]: Requirement Analysis & Domain Research\nT2 [Week 2-3]: UI/UX Design & Frontend Implementation\nT3 [Week 4]: Dataset Generation (Synthetic Modeling)\nT4 [Week 5]: ML Model Training & Optimization\nT5 [Week 6]: System Integration & Delivery Testing"

    # Slide 6: System Flowchart & Architecture
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "System Flowchart & Architecture"
    content = slide.placeholders[1]
    content.text = "Frontend Integration:\nUser Request -> UI Frontend (HTML/CSS)\nUI -> JavaScript Runtime\n\nBackend Integration:\nDataset Generation (generate.py) -> CSV Creation\nML Training Pipeline (train.py) -> Decision Tree Eval (test.py)\n\nResult:\nPrediction -> User"

    # Slide 7: Method of Solution (SRS)
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Method of Solution (SRS)"
    content = slide.placeholders[1]
    content.text = "Input Process:\nUsers input precise climate metrics (Temp/Rain/Humidity).\n\nOutput:\nModel predicts class label (Crop Name) + rich metadata.\n\nData Storage:\nLocal formats (.csv for ML, .js for Catalog) -> scalable.\n\nMenu Design:\nHero Menu, Analysis Tool, Tabulated Results Pane."

    # Slide 8: Programming Code Details
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Programming Code Details"
    content = slide.placeholders[1]
    content.text = "Standardization:\nES6 JavaScript, Modular Python imports, HTML5 layout.\n\nComments & Documentation:\nMath operations & dataset generations thoroughly commented.\n\nError Handling:\nTry/Except in Python.\nMemoryError catching for 32-bit limits.\nUI fallback (Crop Not Found)."

    # Slide 9: Testing Details
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Testing Details"
    content = slide.placeholders[1]
    content.text = "Test Plan:\nValidate dateset integrity, ML accuracy, and workflows.\n\nTest Cases Executed:\nWarm/Wet: Temp 28.5C, Rain 210mm -> Rice (PASS)\nCool/Dry: Temp 18.0C, Rain 80mm -> Wheat (PASS)\n\nTesting Methodologies:\nBlack Box Testing (ML predictions)\nIntegration Testing (raw .csv to Pandas Dataframe)"

    current_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(current_dir, "Smart_Crop_Prediction_Presentation.pptx")
    prs.save(output_path)
    print("Presentation saved to " + output_path)

if __name__ == '__main__':
    main()
