import os
from pptx import Presentation
from pptx.util import Inches, Pt

def main():
    template_path = r"c:\Users\raman\Downloads\College PPT Templet.pptx"
    prs = Presentation(template_path)
    
    xml_slides = prs.slides._sldIdLst
    slides_to_remove = list(xml_slides)[1:]
    for slide in slides_to_remove:
        xml_slides.remove(slide)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Smart Crop Weather Intelligence\n(Weather-Based Crop Prediction)"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Domain: Machine Learning & Agricultural Data Analytics\nAbstract: Analyzing climate inputs to recommend optimal crops using Decision Trees."

    # Slide 2: Intro
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Introduction"
    slide.placeholders[1].text = "Scope & Objectives:\n- Provide accessible, reliable crop recommendations.\n- Maximize yield and mitigate financial risks.\n\nSystem Description:\n- Full-stack app with dynamic web GUI mapped to robust Python ML backend.\n\nProposed System:\n- Responsive search portal with automated pipeline generating 1000+ crops."

    # Slide 3: Overall Description
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Overall Description"
    slide.placeholders[1].text = "Hardware Requirements:\n- Core CPU, 2GB+ RAM, 500MB+ Storage\n\nSoftware Requirements:\n- Windows/Linux/macOS\n- Python, Web Browser\n- Pandas, Scikit-Learn, HTML/JS/CSS\n\nCharacteristics:\n- Responsive UI, scalable DB\n- Cross-browser compatible"

    # Slide 4: Specific Requirements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Specific Requirements"
    slide.placeholders[1].text = "Functional Requirements:\n- Accept temp, rain, humidity inputs -> return prediction.\n- Catalog search for localized needs.\n\nNon-Functional Requirement:\n- Scalable to 1000+ crops, fast load speed.\n\nExternal Interfaces:\n- Particle.js backgrounds, FontAwesome.\n\nPerformance:\n- ML execution avoids memory overflow; Lookups < 500ms."

    # Slide 5: Action Plan
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Action Plan (Gantt Overview)"
    slide.placeholders[1].text = "T1 [Week 1]: Requirement Analysis & Domain Research\nT2 [Week 2-3]: UI/UX Design & Frontend Implementation\nT3 [Week 4]: Dataset Generation (Synthetic Modeling)\nT4 [Week 5]: ML Model Training & Optimization\nT5 [Week 6]: System Integration & Delivery Testing"

    # Slide 6: System Flowchart & Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Flowchart & Architecture"
    slide.placeholders[1].text = "Frontend Integration:\n- User Request -> UI Frontend (HTML/CSS)\n- UI -> JavaScript Runtime\n\nBackend Integration:\n- Dataset Generation (generate.py) -> CSV Creation\n- ML Training Pipeline (train.py) -> Decision Tree Eval (test.py)\n\nResult:\n- Prediction -> User"

    # Slide 7: Website Interface Screenshot
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Website Interface Preview"
    img_path = r"s:\MINI project sample\crop weather\website\ui_screenshot.png"
    if os.path.exists(img_path):
        left = Inches(1.5)
        top = Inches(1.5)
        height = Inches(5)
        slide.shapes.add_picture(img_path, left, top, height=height)
    else:
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "[Website Interface Screenshot missing]"

    # Slide 8: Method of Solution (SRS)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Method of Solution (SRS)"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Input Process:\n- Precise climate metrics (Temp/Rain/Humidity).\n\nOutput:\n- Model predicts class label (Crop Name) + rich metadata.\n\nData Storage:\n- Local formats (.csv for ML, .js for Catalog) -> scalable.\n\nMenu Design:\n- Hero Menu, Analysis Tool, Tabulated Results Pane."

    # Slide 9: Programming Code Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Programming Code Details"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Standardization:\n- ES6 JavaScript, Modular Python imports, HTML5 layout.\n\nComments & Documentation:\n- Math operations & dataset generations thoroughly commented.\n\nError Handling:\n- Try/Except in Python.\n- MemoryError catching for 32-bit limits.\n- UI fallback (Crop Not Found)."

    # Slide 10: Testing Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Testing Details"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Test Plan:\n- Validate dateset integrity, ML accuracy, and workflows.\n\nTest Cases Executed:\n- Warm/Wet: Temp 28.5C, Rain 210mm -> Rice (PASS)\n- Cool/Dry: Temp 18.0C, Rain 80mm -> Wheat (PASS)\n\nTesting Methodologies:\n- Black Box Testing (ML predictions)\n- Integration Testing (raw .csv to Pandas Dataframe)"
    
    out_path = r"s:\MINI project sample\crop weather\website\Final_Project_Presentation.pptx"
    prs.save(out_path)
    print("Successfully built the PPTX at:", out_path)

if __name__ == '__main__':
    main()
